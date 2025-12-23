"""
File conversion utilities for DocMerge.
Handles conversion of various file formats to PDF.
"""

import os
from io import BytesIO
from pathlib import Path

from PIL import Image, ExifTags
from pypdf import PdfReader
from reportlab.lib.pagesizes import A4, LETTER
from reportlab.lib.units import inch
from reportlab.pdfgen import canvas
from reportlab.lib.utils import ImageReader

try:
    from docx import Document
    DOCX_SUPPORT = True
except ImportError:
    DOCX_SUPPORT = False

try:
    from pptx import Presentation
    from pptx.util import Inches
    PPTX_SUPPORT = True
except ImportError:
    PPTX_SUPPORT = False

IMAGE_EXTENSIONS = {'.jpg', '.jpeg', '.png', '.gif', '.webp', '.jfif', '.bmp', '.tiff', '.heic'}
MARKDOWN_EXTENSIONS = {'.md', '.markdown', '.txt'}


def convert_image_to_pdf_bytes(image_path: str, add_source_label: bool = True, page_size=A4) -> bytes | None:
    """
    Convert an image to PDF bytes with optional source label.
    
    Args:
        image_path: Path to the image file
        add_source_label: Whether to add filename label at bottom
        page_size: Page size tuple (width, height)
    
    Returns:
        PDF bytes or None if conversion fails
    """
    try:
        img = Image.open(image_path)
        
        # Convert to RGB if necessary
        if img.mode in ('RGBA', 'P', 'LA'):
            background = Image.new('RGB', img.size, (255, 255, 255))
            if img.mode == 'P':
                img = img.convert('RGBA')
            if img.mode in ('RGBA', 'LA'):
                background.paste(img, mask=img.split()[-1])
            img = background
        elif img.mode != 'RGB':
            img = img.convert('RGB')
        
        # Auto-rotate based on EXIF orientation
        try:
            exif = img._getexif()
            if exif:
                for tag, value in exif.items():
                    if ExifTags.TAGS.get(tag) == 'Orientation':
                        if value == 3:
                            img = img.rotate(180, expand=True)
                        elif value == 6:
                            img = img.rotate(270, expand=True)
                        elif value == 8:
                            img = img.rotate(90, expand=True)
                        break
        except (AttributeError, KeyError, TypeError):
            pass
        
        # Calculate scaling
        img_width, img_height = img.size
        page_width, page_height = page_size
        
        margin = 0.5 * inch
        available_width = page_width - 2 * margin
        available_height = page_height - 2 * margin - (0.5 * inch if add_source_label else 0)
        
        scale_w = available_width / img_width
        scale_h = available_height / img_height
        scale = min(scale_w, scale_h)
        
        new_width = img_width * scale
        new_height = img_height * scale
        
        # Center image
        x = (page_width - new_width) / 2
        y = (page_height - new_height) / 2
        if add_source_label:
            y += 0.25 * inch
        
        # Create PDF
        buffer = BytesIO()
        c = canvas.Canvas(buffer, pagesize=page_size)
        
        temp_img = BytesIO()
        img.save(temp_img, format='PNG')
        temp_img.seek(0)
        
        c.drawImage(ImageReader(temp_img), x, y, width=new_width, height=new_height)
        
        if add_source_label:
            filename = os.path.basename(image_path)
            c.setFont("Helvetica", 8)
            c.drawString(margin, margin, f"Source: {filename}")
        
        c.save()
        buffer.seek(0)
        return buffer.getvalue()
    except Exception as e:
        print(f"  ERROR converting image {image_path}: {e}")
        return None


def convert_docx_to_pdf_bytes(docx_path: str, page_size=A4) -> bytes | None:
    """
    Convert DOCX to PDF bytes using text extraction.
    
    Args:
        docx_path: Path to the DOCX file
        page_size: Page size tuple
    
    Returns:
        PDF bytes or None if conversion fails
    """
    if not DOCX_SUPPORT:
        print(f"  WARNING: python-docx not installed, skipping {docx_path}")
        return None
    
    try:
        doc = Document(docx_path)
        buffer = BytesIO()
        c = canvas.Canvas(buffer, pagesize=page_size)
        
        page_width, page_height = page_size
        margin = inch
        y = page_height - margin
        line_height = 14
        max_width = page_width - 2 * margin
        
        # Document title
        filename = os.path.basename(docx_path)
        c.setFont("Helvetica-Bold", 12)
        c.drawString(margin, y, f"Document: {filename}")
        y -= line_height * 2
        
        c.setFont("Helvetica", 10)
        
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                y -= line_height / 2
                continue
            
            # Word wrapping
            words = text.split()
            lines = []
            current_line = []
            
            for word in words:
                test_line = ' '.join(current_line + [word])
                if c.stringWidth(test_line, "Helvetica", 10) < max_width:
                    current_line.append(word)
                else:
                    if current_line:
                        lines.append(' '.join(current_line))
                    current_line = [word]
            if current_line:
                lines.append(' '.join(current_line))
            
            for line in lines:
                if y < margin:
                    c.showPage()
                    c.setFont("Helvetica", 10)
                    y = page_height - margin
                c.drawString(margin, y, line)
                y -= line_height
            
            y -= line_height / 2
        
        # Extract tables
        for table in doc.tables:
            y -= line_height
            if y < margin * 2:
                c.showPage()
                c.setFont("Helvetica", 10)
                y = page_height - margin
            
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if y < margin:
                    c.showPage()
                    c.setFont("Helvetica", 10)
                    y = page_height - margin
                while c.stringWidth(row_text, "Helvetica", 10) > max_width and len(row_text) > 10:
                    row_text = row_text[:-10] + "..."
                c.drawString(margin, y, row_text)
                y -= line_height
        
        c.save()
        buffer.seek(0)
        return buffer.getvalue()
    except Exception as e:
        print(f"  ERROR converting DOCX {docx_path}: {e}")
        return None


def create_title_page(title: str, subtitle: str = "", page_size=A4) -> bytes:
    """
    Create a title page PDF.
    
    Args:
        title: Main title text
        subtitle: Optional subtitle
        page_size: Page size tuple
    
    Returns:
        PDF bytes
    """
    buffer = BytesIO()
    c = canvas.Canvas(buffer, pagesize=page_size)
    page_width, page_height = page_size
    
    c.setFont("Helvetica-Bold", 24)
    c.drawCentredString(page_width / 2, page_height - 3 * inch, title)
    
    if subtitle:
        c.setFont("Helvetica", 14)
        c.drawCentredString(page_width / 2, page_height - 3.5 * inch, subtitle)
    
    c.setFont("Helvetica", 12)
    c.drawCentredString(page_width / 2, page_height - 5 * inch, "Generated by DocMerge")
    
    c.save()
    buffer.seek(0)
    return buffer.getvalue()


def convert_markdown_to_pdf_bytes(md_path: str, page_size=A4) -> bytes | None:
    """Convert markdown/text file to PDF bytes."""
    try:
        with open(md_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        
        buffer = BytesIO()
        c = canvas.Canvas(buffer, pagesize=page_size)
        page_width, page_height = page_size
        margin = inch
        y = page_height - margin
        line_height = 12
        max_width = page_width - 2 * margin
        
        filename = os.path.basename(md_path)
        c.setFont("Helvetica-Bold", 12)
        c.drawString(margin, y, f"Document: {filename}")
        y -= line_height * 2
        
        c.setFont("Courier", 9)
        
        for line in content.split('\n'):
            if y < margin:
                c.showPage()
                c.setFont("Courier", 9)
                y = page_height - margin
            
            # Handle long lines
            while len(line) > 0:
                if c.stringWidth(line, "Courier", 9) <= max_width:
                    c.drawString(margin, y, line)
                    y -= line_height
                    break
                else:
                    # Find break point
                    break_at = len(line)
                    while break_at > 0 and c.stringWidth(line[:break_at], "Courier", 9) > max_width:
                        break_at -= 1
                    if break_at == 0:
                        break_at = 1
                    c.drawString(margin, y, line[:break_at])
                    y -= line_height
                    line = line[break_at:]
                    if y < margin:
                        c.showPage()
                        c.setFont("Courier", 9)
                        y = page_height - margin
        
        c.save()
        buffer.seek(0)
        return buffer.getvalue()
    except Exception as e:
        print(f"  ERROR converting markdown {md_path}: {e}")
        return None


def convert_pptx_to_pdf_bytes(pptx_path: str, page_size=A4) -> bytes | None:
    """Convert PowerPoint file to PDF bytes."""
    if not PPTX_SUPPORT:
        print("  WARNING: python-pptx not installed, cannot convert PPTX")
        return None
    
    try:
        prs = Presentation(pptx_path)
        buffer = BytesIO()
        c = canvas.Canvas(buffer, pagesize=page_size)
        page_width, page_height = page_size
        margin = 0.5 * inch
        
        filename = os.path.basename(pptx_path)
        
        for slide_num, slide in enumerate(prs.slides, 1):
            # Title for each slide
            c.setFont("Helvetica-Bold", 14)
            c.drawString(margin, page_height - margin, f"{filename} - Slide {slide_num}")
            
            y = page_height - margin - 30
            c.setFont("Helvetica", 11)
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    for line in text.split('\n'):
                        if y < margin + 20:
                            c.showPage()
                            c.setFont("Helvetica-Bold", 14)
                            c.drawString(margin, page_height - margin, f"{filename} - Slide {slide_num} (cont.)")
                            y = page_height - margin - 30
                            c.setFont("Helvetica", 11)
                        
                        # Truncate long lines
                        max_chars = int((page_width - 2 * margin) / 6)
                        if len(line) > max_chars:
                            line = line[:max_chars-3] + "..."
                        c.drawString(margin, y, line)
                        y -= 14
            
            c.showPage()
        
        c.save()
        buffer.seek(0)
        return buffer.getvalue()
    except Exception as e:
        print(f"  ERROR converting PPTX {pptx_path}: {e}")
        return None


def is_supported_file(filepath: str) -> bool:
    """Check if file type is supported for conversion."""
    ext = Path(filepath).suffix.lower()
    return (ext in IMAGE_EXTENSIONS or ext == '.pdf' or ext == '.docx' 
            or ext in MARKDOWN_EXTENSIONS or ext == '.pptx')


def get_file_type(filepath: str) -> str:
    """Get the type category of a file."""
    ext = Path(filepath).suffix.lower()
    if ext == '.pdf':
        return 'pdf'
    elif ext == '.docx':
        return 'docx'
    elif ext == '.pptx':
        return 'pptx'
    elif ext in IMAGE_EXTENSIONS:
        return 'image'
    elif ext in MARKDOWN_EXTENSIONS:
        return 'markdown'
    return 'unknown'
